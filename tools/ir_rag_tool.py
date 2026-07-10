"""
tools/ir_rag_tool.py
RAG pipeline that automatically finds, downloads and analyses IR documents.

Required packages (in addition to existing requirements):
    pip install langchain-community pypdf faiss-cpu tiktoken python-pptx beautifulsoup4
"""

import os
import json
import re
import time
import hashlib
from pathlib import Path
from urllib.parse import urljoin, urlparse

import requests
import yfinance as yf
from bs4 import BeautifulSoup
from dotenv import load_dotenv
from langchain_core.tools import tool
from langchain_core.documents import Document
from langchain_openai import OpenAIEmbeddings
from langchain_anthropic import ChatAnthropic
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.output_parsers import StrOutputParser

try:
    from langchain_text_splitters import RecursiveCharacterTextSplitter
except ImportError:
    from langchain.text_splitter import RecursiveCharacterTextSplitter  # type: ignore

from langchain_community.document_loaders import PyPDFLoader
from langchain_community.vectorstores import FAISS

try:
    from langchain_community.document_loaders import Docx2txtLoader
    _DOCX_AVAILABLE = True
except ImportError:
    _DOCX_AVAILABLE = False

try:
    from pptx import Presentation as _PptxPresentation
    _PPTX_AVAILABLE = True
except ImportError:
    _PPTX_AVAILABLE = False

try:
    from playwright.sync_api import sync_playwright as _sync_playwright
    _PLAYWRIGHT_AVAILABLE = True
except ImportError:
    _PLAYWRIGHT_AVAILABLE = False

import sys
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from tools.finance_tools import get_cashflow_data

load_dotenv()

# ── Constants ─────────────────────────────────────────────────────────────────

IR_PAGE_PATTERNS = [
    # Bestehende Muster
    "/investors/publications",
    "/investors",
    "/investor-relations",
    "/ir",
    "/investors/financial-results",
    "/investors/financials/financial-publications",
    "/financials/annual-reports",
    "/financial-info/annual-reports",
    "/annual-reporting-suite",
    "/investors/results-reports-presentations",
    "/annual-report",
    "/financial-performance",
    "/reporting-center",
    "/investors/financial-data",
    "/shareholders",
    "/investors/reports",
    "/investors/reports-publications",
    "/results-and-reports",
    "/financial-information",
    "/results-and-presentations",
    "/investors/financial-reporting",
    "/media-investors",
    "/results-centre",
    "/company/investors",
    "/about/investors",
    "/about-us/investors",
    "/en/investors",
    "/de/investoren",
    "/group/investors",
    "/corporate/investors",
    "/investisseurs",
    "/rapports-annuels",
    "/resultats-financiers",
    "/investoren",
    "/investor_relations",
    "/berichte-und-praesentationen",
    "/finanzberichte",
    "/geschaeftsberichte",
    "/publikationen",
    "/finanzpublikationen",
    "/ergebnisse",
    "/download-center",
    "/investor-relations/downloads",
    "/investor-relations/reports-and-presentations",
    "/investor-relations/financial-reports",
    "/investor-relations/news",
    "/financial-reports",
    "/annual-results",
    "/reporting",
    "/reports",
    "/downloads",
    "/archive",
    "/investor-relations-downloads",
    "/investors-hub",
    "/media/publications",
    "/en/group/investor-relations",
]

PRIORITY_KEYWORDS = [
    "consensus",
    "analyst presentation",
    "annual report",
    "financial report",
    "half-year",
    "quarterly results",
    "investor day",
    "outlook",
    "guidance",
    "full year",
    "integrated report",
    "key figures",
    "financial statements",
    "factsheet",
    "presentation",
    "earnings release",
    "letter to shareholders",
    "investor presentation",
    "interim report",
    "geschäftsbericht",
    "rapport annuel",
]

EXCLUDE_KEYWORDS = [
    "sustainability",
    "governance",
    "compensation",
    "esg",
    "proxy",
]

CACHE_DIR = "./ir_cache"
CACHE_MAX_AGE_HOURS = 24

_KNOWN_IR_URLS: dict[str, str] = {
    # Swiss
    "ABBN.SW": "https://www.abb.com/global/en/company/annual-reporting-suite#download",
    "ALC.SW":  "https://investor.alcon.com/financials/annual-reports",
    "AMRZ.SW": "https://investors.amrize.com/financial-info/annual-reports",
    "CFR.SW":  "https://www.richemont.com/investors/results-reports-presentations/",
    "GEBN.SW": "https://reports.geberit.com/annual-report/2025/services/downloads.html#downloads-archive",
    "GIVN.SW": "https://www.givaudan.com/investors/financial-results/results-centre",
    "HOLN.SW": "https://www.holcim.com/investors/publications",
    "KNIN.SW": "https://www.kuehne-nagel.com/company/investor-relations/financial-performance",
    "LOGN.SW": "https://ir.logitech.com/financial-info/annual-reports/default.aspx",
    "LONN.SW": "https://www.lonza.com/investor-relations/reporting-center",
    "NESN.SW": "https://www.nestle.com/investors/publications",
    "NOVN.SW": "https://www.novartis.com/investors/financial-data/annual-results",
    "PGHN.SW": "https://www.partnersgroup.com/en/shareholders/reports-and-presentations#all",
    "ROP.SW":  "https://www.roche.com/investors/reports#9e6fe792-f417-4188-993e-9cedf91bd4f6",
    "SIKA.SW": "https://www.sika.com/en/investors/reports-publications/financial-reports.html",
    "SLHN.SW": "https://www.swisslife.com/en/home/investors/results-and-reports.html",
    "SREN.SW": "https://www.swissre.com/investors/financial-information.html#2025-content",
    "SCMN.SW": "https://www.swisscom.ch/en/about/investors/reports.html",
    "UBSG.SW": "https://www.ubs.com/global/en/investor-relations/financial-information/sec-filings.html#tab-1955586228",
    "ZURN.SW": "https://www.zurich.com/investor-relations/results-and-reports",
    "RIEN.SW": "https://www.rieter.com/investor-relations/results-and-presentations/financial-reports",
    "SCHN.SW": "https://www.schindler.com/com/internet/en/investor-relations/reports.html",
    "BAER.SW": "https://www.juliusbaer.com/en/media-investors/financial-information/financial-reporting-1/#c101123",
    "LISN.SW": "https://www.lindt-spruengli.com/investors/financial-reporting/publications",
    "AAPL":    "https://investor.apple.com/financial-information/sec-filings/default.aspx",
    "MSFT":    "https://www.microsoft.com/en-us/investor/sec-filings.aspx",
    "GOOGL":   "https://abc.xyz/investor/",
    "GOOG":    "https://abc.xyz/investor/",
    "META":    "https://investor.atmeta.com/sec-filings/annual-reports",
    "AMZN":    "https://ir.aboutamazon.com/sec-filings/annual-reports",
    "NVDA":    "https://investor.nvidia.com/financial-information/sec-filings",
    "TSLA":    "https://ir.tesla.com/sec-filings/annual-reports",
    "JPM":     "https://www.jpmorganchase.com/ir/annual-report",
    "GS":      "https://www.goldmansachs.com/investor-relations/financials/",
}

# Foreign private issuers: maps exchange ticker -> US SEC ticker for EDGAR lookup
# These companies file 20-F (annual) and 6-K (interim) instead of 10-K/10-Q
_SEC_FOREIGN_FILERS: dict[str, str] = {
    "UBSG.SW": "UBS",
    "NOVN.SW": "NVS",
    "ABBN.SW": "ABBNY",  # ABB Ltd ADR — files 20-F/6-K with SEC
}

# HTML anchor-text keywords used to detect IR documents when no PDFs are found
_HTML_IR_KEYWORDS = [
    "annual report", "quarterly results", "earnings", "press release",
    "presentation", "outlook", "guidance", "10-k", "10-q", "10k", "10q",
    "investor presentation", "financial results", "interim report",
]

# URL path segments that indicate a report-index sub-page (EU/CH multi-level IR sites)
_SUBPAGE_FOLLOW_PATTERNS = [
    "financial-report", "annual-report", "half-year-report", "interim-report",
    "results-and-presentations", "financial-results", "results",
    "publications", "downloads", "filings", "reports",
    "geschaeftsbericht", "jahresbericht", "halbjahresbericht", "berichte",
    "rapport-annuel", "resultats", "publications-financieres",
]

# EU/CH PDF fallback classification when _PDF_TYPE_RULES don't match
_EU_PDF_ANNUAL  = re.compile(
    r"annual|full[_\-\s]year|geschaeft|jahres|rapport[_\-\s]annuel|full[_\-\s]report",
    re.I,
)
_EU_PDF_INTERIM = re.compile(
    r"half[_\-\s]year|interim|halbjahr|semest|six[_\-\s]month|h1[-_\s\.]|h2[-_\s\.]",
    re.I,
)

# SEC EDGAR base URLs and constants
SEC_USER_AGENT         = "KI-Co-Analyst research@bfh.ch"
SEC_TICKERS_CACHE      = Path(CACHE_DIR) / "sec_tickers.json"
SEC_TICKERS_URL        = "https://www.sec.gov/files/company_tickers.json"
SEC_TICKERS_MAX_AGE_DAYS = 7

_SEC_EDGAR_BASE      = "https://www.sec.gov"
_SEC_HEADERS         = {
    "User-Agent": SEC_USER_AGENT,
    "Accept":     "application/json",
}
_SEC_TICKERS_CACHE   = SEC_TICKERS_CACHE
_SEC_TICKERS_CACHE_H = SEC_TICKERS_MAX_AGE_DAYS * 24

_HEADERS = {
    "User-Agent": (
        "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
        "AppleWebKit/537.36 (KHTML, like Gecko) "
        "Chrome/120.0.0.0 Safari/537.36"
    ),
    "Accept": "text/html,application/xhtml+xml,*/*;q=0.8",
}

_MAX_FILE_BYTES = 50 * 1024 * 1024  # 50 MB

# Doc type → (priority, keyword pairs that must ALL appear in label)
_PDF_TYPE_RULES: list[tuple[str, int, list[str], list[str]]] = [
    # (type, priority, required_any, required_all)
    ("consensus_estimates",  1, ["consensus"],                      []),
    ("analyst_presentation", 2, ["analyst"],                        ["presentation"]),
    ("annual_report",        3, ["annual report", "financial report"], []),
    ("annual_report",        3, ["annual"],                         ["report"]),
    ("annual_report",        3, ["geschaeftsbericht", "jahresbericht", "rapport annuel"], []),
    ("interim_report",       4, ["half-year", "half year", "h1 results", "h2 results",
                                 "interim", "halbjahr", "six month"],  []),
    ("investor_day",         5, ["investor day", "investor presentation", "cmd"], []),
]

STANDARD_QUERIES = [
    # 1. Die "Goldgrube": Übersichts-Tabellen & Key Figures
    "Financial highlights key figures summary table five-year overview",
    "Kennzahlen Konzernrechnung Gruppenübersicht Mehrjahresvergleich",
    
    # 2. Erfolgsrechnung (Für EBITDA, EBIT, Profit)
    "Consolidated Income Statement Earnings before interest taxes depreciation EBITDA EBIT",
    "Konsolidierte Erfolgsrechnung Betriebsergebnis vor Abschreibungen",
    
    # 3. Bilanz (Für ROE, ROIC, Assets)
    "Consolidated Balance Sheet total assets equity liabilities cash debt",
    "Konsolidierte Bilanz Bilanzsumme Eigenkapital Nettoverschuldung",
    
    # 4. Cashflow (Für FCF & Investitionen)
    "Consolidated Statement of Cash Flows operating investing free cash flow FCF",
    "Geldflussrechnung Investitionen Sachanlagen Capex",
    
    # 5. EPS & Guidance (Für Bewertung & Ausblick)
    "Adjusted EPS earnings per share bereinigter Gewinn pro Aktie restated",
    "Guidance outlook forecast target profit margin Umsatzprognose Ausblick",
    
    # 6. Marktdaten & Konsens
    "Consensus estimates analyst forecast revenue EPS 2026 2027 2028",
    "Dividend policy payout per share Dividende Ausschüttung",
    
    # 7. Sektor-spezifisch (Für Industrials wie Rieter)
    "Order intake order backlog Auftragseingang Auftragsbestand",
]

_splitter = RecursiveCharacterTextSplitter(chunk_size=1500, chunk_overlap=200)


# ── Lazy-init LLM / Embeddings ────────────────────────────────────────────────

_llm: ChatAnthropic | None = None
_emb: OpenAIEmbeddings | None = None


def _get_llm() -> ChatAnthropic:
    global _llm
    if _llm is None:
        _llm = ChatAnthropic(model="claude-sonnet-4-6", temperature=0)
    return _llm


def _get_emb() -> OpenAIEmbeddings:
    global _emb
    if _emb is None:
        _emb = OpenAIEmbeddings()
    return _emb


# ── Document deduplication & year spreading ───────────────────────────────────

def _deduplicate_and_spread(docs: list[dict],
                            max_annual: int = 3,
                            max_latest: int = 4,
                            wanted_years: set[int] | None = None) -> list[dict]:
    """
    1. Remove language duplicates (EN > DE > FR > IT preferred).
    2. Stamp period_class ("annual" | "quarterly") for EU/CH docs that lack it.
    3. Select annual docs + up to max_latest most-recent non-annual docs from
       the current/prior fiscal year (quarterly/interim preferred), so every
       interim report published so far this year (Q1, H1, 9M, ...) can be
       captured.

    If *wanted_years* is given, annual docs are picked to cover exactly those
    fiscal years (DB gaps + latest year for restatement checks) instead of
    just the newest max_annual years — this lets a mature ticker top up a
    single missing year without re-extracting years already cached.

    Returns annual_docs + latest_docs.
    """
    from datetime import datetime as _dt
    from urllib.parse import unquote as _unquote
    from tools.period_classifier import classify_pdf_period

    if not docs:
        return docs

    def _lang_rank(doc: dict) -> int:
        combined = (_unquote(doc.get("url", "")) + " " + doc.get("filename", "")).lower()
        for suffix, rank in (("_en.", 0), ("/en/", 0), ("_de.", 1), ("/de/", 1),
                              ("_fr.", 2), ("_it.", 3)):
            if suffix in combined:
                return rank
        return 2

    def _norm_url(url: str) -> str:
        decoded = _unquote(url).lower()
        return re.sub(r"[_-](?:de|en|fr|it)(?=\.pdf)", "", decoded)

    def _infer_period_class(doc: dict) -> str:
        """Infer period_class for EU/CH docs that don't carry it from get_sec_filings()."""
        hint = (doc.get("filename", "") + " " + doc.get("text", "") +
                " " + doc.get("type", "")).lower()
        pt = classify_pdf_period(hint)
        if pt == "annual":
            return "annual"
        if pt in ("quarterly", "h1", "9m"):
            return "quarterly"
        # Fallback: use the existing 'type' field
        return "annual" if doc.get("type") == "annual_report" else "quarterly"

    # Stamp period_class if missing
    for doc in docs:
        if "period_class" not in doc:
            doc["period_class"] = _infer_period_class(doc)

    # Group by (period_class, year, normalised_url) → keep best language per group
    groups: dict[tuple, dict] = {}
    for doc in docs:
        key = (doc["period_class"], doc["year"],
               _norm_url(doc.get("url", doc["filename"])))
        if key not in groups or _lang_rank(doc) < _lang_rank(groups[key]):
            groups[key] = doc

    current_year = _dt.now().year
    min_year     = current_year - (max(current_year - min(wanted_years), 6) if wanted_years else 6)

    annual_candidates = sorted(
        [d for d in groups.values() if d["period_class"] == "annual"
         and (d["year"] == 0 or d["year"] >= min_year)],
        key=lambda d: -d["year"],
    )
    other_candidates = sorted(
        [d for d in groups.values() if d["period_class"] != "annual"
         and (d["year"] == 0 or d["year"] >= min_year)],
        key=lambda d: (d["priority"] == 2 and d["period_class"] == "quarterly", -d["year"]),
        reverse=True,
    )

    if wanted_years:
        # Gap-driven: keep only annual docs whose year is actually missing
        # from the DB (plus year==0 docs, since their real year is unknown
        # until parsed — better to include than silently drop them).
        annual_result = [d for d in annual_candidates
                          if d["year"] == 0 or d["year"] in wanted_years]
    else:
        # Pick up to max_annual distinct years for annuals
        annual_result = []
        seen_annual_years: set = set()
        for doc in annual_candidates:
            yr = doc["year"] if doc["year"] > 0 else f"_u{len(annual_result)}"
            if yr not in seen_annual_years:
                seen_annual_years.add(yr)
                annual_result.append(doc)
            if len(annual_result) >= max_annual:
                break

    # Pick up to max_latest interim docs from the current/prior fiscal year
    # (prefer quarterly), so all interim reports published so far this year
    # are captured — not just the single newest one.
    recent_other = [d for d in other_candidates if d["year"] == 0 or d["year"] >= current_year - 1]
    quarterly_first = sorted(
        recent_other,
        key=lambda d: (0 if d["period_class"] == "quarterly" else 1, -d.get("year", 0)),
    )
    latest_result = quarterly_first[:max_latest]

    return annual_result + latest_result


# ── Playwright helper ─────────────────────────────────────────────────────────

def _playwright_render(url: str, timeout_ms: int = 20_000) -> str:
    """
    Renders *url* with a headless Chromium and returns the full HTML after JS execution.
    Returns '' if Playwright is not installed or on any error.
    """
    if not _PLAYWRIGHT_AVAILABLE:
        return ""
    try:
        with _sync_playwright() as pw:
            browser = pw.chromium.launch(headless=True)
            ctx  = browser.new_context(user_agent=_HEADERS["User-Agent"])
            page = ctx.new_page()
            page.goto(url, timeout=timeout_ms, wait_until="commit")
            try:
                page.wait_for_load_state("networkidle", timeout=10_000)
            except Exception:
                page.wait_for_timeout(3_000)
            try:
                page.evaluate("""
                    () => {
                        const patterns = [
                            '#onetrust-accept-btn-handler',
                            '#accept-recommended-btn-handler',
                            '#CybotCookiebotDialogBodyButtonAccept',
                            'button[id*="accept"]','button[class*="accept"]',
                        ];
                        for (const sel of patterns) {
                            const btn = document.querySelector(sel);
                            if (btn) { btn.click(); return true; }
                        }
                        const texts = ['accept all','accept all cookies','alle akzeptieren',
                                       'akzeptieren','i accept','agree','ok'];
                        for (const btn of document.querySelectorAll('button')) {
                            if (texts.includes(btn.textContent.trim().toLowerCase())) {
                                btn.click(); return true;
                            }
                        }
                        return false;
                    }
                """)
                page.wait_for_timeout(1_500)
            except Exception:
                pass
            html = page.content()
            browser.close()
            return html
    except Exception as exc:
        print(f"      Playwright Fehler ({url[:60]}): {exc}")
        return ""


def _playwright_render_multi_year(url: str, n_years: int = 5,
                                   timeout_ms: int = 20_000) -> str:
    """
    Renders *url* and then clicks year-filter buttons for the n most recent years.
    Many EU/CH IR pages (Swiss Life, Zurich, etc.) show only the current year by
    default and require clicking year-tabs to reveal older documents.

    Returns concatenated HTML from all year views so PDF scraping covers all years.
    Falls back silently to a single-page render if no year buttons are found.
    """
    if not _PLAYWRIGHT_AVAILABLE:
        return ""

    from datetime import datetime as _dt
    current_year = _dt.now().year
    # Try previous years (current year is visible on initial load)
    prev_years = list(range(current_year - 1, current_year - n_years - 1, -1))

    combined_html = ""
    try:
        with _sync_playwright() as pw:
            browser = pw.chromium.launch(headless=True)
            ctx  = browser.new_context(
                user_agent=_HEADERS["User-Agent"],
                viewport={"width": 1920, "height": 1080},
            )
            page = ctx.new_page()
            # 'commit' fires as soon as response headers arrive — never times out on slow
            # sites that delay DOMContentLoaded (e.g. Zurich Insurance). Then we wait
            # for networkidle (up to 20 s) or fall back to a fixed 5-second pause.
            page.goto(url, timeout=45_000, wait_until="commit")
            try:
                page.wait_for_load_state("networkidle", timeout=20_000)
            except Exception:
                page.wait_for_timeout(5_000)

            # Accept cookie consent banners (OneTrust, Cookiebot, generic)
            try:
                page.evaluate("""
                    () => {
                        const patterns = [
                            '#onetrust-accept-btn-handler',
                            '#accept-recommended-btn-handler',
                            '#CybotCookiebotDialogBodyButtonAccept',
                            'button[id*="accept"]',
                            'button[class*="accept"]',
                        ];
                        for (const sel of patterns) {
                            const btn = document.querySelector(sel);
                            if (btn) { btn.click(); return true; }
                        }
                        // Text-based fallback
                        const texts = ['accept all','accept all cookies','alle akzeptieren',
                                       'akzeptieren','i accept','agree','ok'];
                        for (const btn of document.querySelectorAll('button')) {
                            if (texts.includes(btn.textContent.trim().toLowerCase())) {
                                btn.click(); return true;
                            }
                        }
                        return false;
                    }
                """)
                page.wait_for_timeout(1_500)
            except Exception:
                pass

            combined_html += page.content()

            for year in prev_years:
                try:
                    clicked = page.evaluate(f"""
                        () => {{
                            const btn = Array.from(document.querySelectorAll("button"))
                                .find(b => b.textContent.trim() === "{year}");
                            if (btn) {{ btn.click(); return true; }}
                            return false;
                        }}
                    """)
                    if clicked:
                        page.wait_for_timeout(1_500)
                        combined_html += page.content()
                        print(f"      Playwright Tab {year}: OK")
                except Exception:
                    continue

            browser.close()
    except Exception as exc:
        print(f"      Playwright multi-year Fehler ({url[:60]}): {exc}")

    return combined_html


# ── 1. Find IR URL ────────────────────────────────────────────────────────────

def find_ir_url(ticker: str, company_name: str = "") -> str:
    """
    Returns the best IR page URL for *ticker*, or '' if not found.

    For US tickers (no '.' in symbol): SEC EDGAR is the primary source.
      Returns a canonical EDGAR browse URL that find_ir_pdfs() detects and
      routes directly to get_sec_filings().

    For EU/CH tickers: tries the known-URL mapping, then yfinance + pattern probing.
    """
    # ── US tickers: SEC EDGAR as primary path ────────────────────────────────
    if _sec_is_us_ticker(ticker):
        cik = get_sec_cik(ticker)
        if cik:
            cik_padded = cik.zfill(10)
            sec_url = (
                f"{_SEC_EDGAR_BASE}/cgi-bin/browse-edgar"
                f"?action=getcompany&CIK={cik_padded}"
                f"&type=10-K&dateb=&owner=include&count=10"
            )
            print(f"      SEC EDGAR: CIK={int(cik)} fuer {ticker} -> {sec_url[:70]}")
            return sec_url
        # CIK not found — fall through to website scraping below

    # ── Known hard-coded mapping (EU/CH primary, US fallback) ────────────────
    # Trusted URLs are returned directly — no HTTP check needed.
    # Many IR sites return 403 to bots (requests) but load fine via Playwright.
    if ticker in _KNOWN_IR_URLS:
        url = _KNOWN_IR_URLS[ticker]
        print(f"      IR-Seite (bekannte URL): {url}")
        return url

    # ── yfinance website + IR path patterns (EU/CH tickers) ──────────────────
    base_url = ""
    try:
        info = yf.Ticker(ticker).info
        website = info.get("website", "").strip().rstrip("/")
        if website:
            if not website.startswith("http"):
                website = "https://" + website
            base_url = website
    except Exception:
        pass

    if not base_url:
        return ""

    for path in IR_PAGE_PATTERNS:
        url = base_url + path
        try:
            r = requests.get(url, headers=_HEADERS, timeout=8, allow_redirects=True)
            if r.status_code == 200 and len(r.content) > 800:
                print(f"      IR-Seite gefunden: {url}")
                return url
        except Exception:
            continue

    return base_url


# ── 2a. SEC EDGAR helpers ─────────────────────────────────────────────────────

def _sec_is_us_ticker(ticker: str) -> bool:
    """Returns True for plain US tickers (no exchange suffix like .SW, .L, .PA)."""
    return "." not in ticker


def get_sec_cik(ticker: str) -> str | None:
    """
    Returns the 10-digit zero-padded CIK for *ticker*, or None.

    Way A (primary): SEC company_tickers.json — full list, cached 7 days locally.
    Way B (fallback): browse-edgar atom feed — parses <cik-number> from XML.
    """
    ticker_upper = ticker.upper()

    # Way A: local cache of the full tickers list (7-day TTL)
    mapping: dict = {}
    cache_ok = False
    try:
        if _SEC_TICKERS_CACHE.exists():
            age_h = (time.time() - _SEC_TICKERS_CACHE.stat().st_mtime) / 3600
            if age_h < _SEC_TICKERS_CACHE_H:
                mapping  = json.loads(_SEC_TICKERS_CACHE.read_text(encoding="utf-8"))
                cache_ok = True
    except Exception:
        pass

    if not cache_ok:
        try:
            r = requests.get(
                SEC_TICKERS_URL,
                headers=_SEC_HEADERS, timeout=12,
            )
            r.raise_for_status()
            mapping = r.json()
            _SEC_TICKERS_CACHE.parent.mkdir(parents=True, exist_ok=True)
            _SEC_TICKERS_CACHE.write_text(
                json.dumps(mapping, ensure_ascii=False), encoding="utf-8"
            )
            time.sleep(0.15)  # SEC rate-limit courtesy
        except Exception:
            pass

    for entry in mapping.values():
        if str(entry.get("ticker", "")).upper() == ticker_upper:
            return str(entry["cik_str"]).zfill(10)

    # Way B: EDGAR atom feed — accepts ticker directly as CIK parameter
    try:
        r = requests.get(
            f"{_SEC_EDGAR_BASE}/cgi-bin/browse-edgar",
            params={
                "company": "", "CIK": ticker, "type": "10-K",
                "dateb": "", "owner": "include", "count": "1",
                "search_text": "", "action": "getcompany", "output": "atom",
            },
            headers=_SEC_HEADERS, timeout=12,
        )
        r.raise_for_status()
        time.sleep(0.15)
        # <cik-number> tag in the atom response
        m = re.search(r"<cik-number>(\d+)</cik-number>", r.text)
        if m:
            return m.group(1).zfill(10)
        # Fallback: 10-digit CIK embedded in a URL inside the response
        m = re.search(r"CIK=(\d{10})", r.text)
        if m:
            return m.group(1)
    except Exception:
        pass

    return None


def _sec_fetch_filing_doc_url(filing_index_url: str) -> str | None:
    """
    Parses the SEC filing index page and returns the primary document URL.
    Used when the submissions JSON doesn't supply a primaryDocument name.
    """
    try:
        r = requests.get(filing_index_url, headers=_SEC_HEADERS, timeout=15)
        r.raise_for_status()
        time.sleep(0.15)
        soup = BeautifulSoup(r.text, "html.parser")
        for row in soup.select("table tr"):
            cells = row.find_all("td")
            if len(cells) < 3:
                continue
            doc_type = cells[0].get_text(strip=True).upper()
            doc_href = cells[2].find("a", href=True) if len(cells) > 2 else None
            if not doc_href:
                continue
            href = doc_href["href"]
            if not href.lower().endswith((".htm", ".html")):
                continue
            if "index" in href.lower():
                continue
            if doc_type in ("10-K", "10-Q", "8-K", "10-K/A", "10-Q/A"):
                return href if href.startswith("http") else _SEC_EDGAR_BASE + href
        # Fallback: first non-index htm link in the table
        for a in soup.select("table a[href]"):
            href = a["href"]
            if href.lower().endswith((".htm", ".html")) and "index" not in href.lower():
                return href if href.startswith("http") else _SEC_EDGAR_BASE + href
    except Exception:
        pass
    return None


def get_sec_filings(ticker: str, cik: str,
                    max_annuals: int = 3,
                    max_quarterly: int = 1) -> list[dict]:
    """
    Fetches up to max_annuals annual filings + max_quarterly quarterly filings
    for *ticker* from SEC EDGAR.

    Annual forms  : 10-K, 20-F, 10-K/A, 20-F/A
    Quarterly forms: 10-Q, 6-K, 10-Q/A

    Each returned dict is compatible with the find_ir_pdfs schema:
      {url, filename, type, priority, year, text, format, source, period_class}

    period_class is "annual" or "quarterly" — used downstream for routing.
    """
    cik_padded = cik.zfill(10)
    cik_int    = int(cik)

    _ANNUAL_FORMS    = {"10-K", "20-F", "10-K/A", "20-F/A"}
    _QUARTERLY_FORMS = {"10-Q", "6-K", "10-Q/A"}
    _PRIORITY = {
        "10-K": 3, "20-F": 3, "10-K/A": 3, "20-F/A": 3,
        "10-Q": 2, "6-K":  2, "10-Q/A": 2,
    }

    annuals:          list[dict] = []
    quarterly_latest: list[dict] = []

    def _build_entry(form_clean: str, acc: str, date: str,
                     prim_doc: str, period_class: str) -> dict:
        acc_nodash   = acc.replace("-", "")
        archive_base = f"{_SEC_EDGAR_BASE}/Archives/edgar/data/{cik_int}/{acc_nodash}"
        doc_url = (
            f"{archive_base}/{prim_doc}"
            if prim_doc
            else _sec_fetch_filing_doc_url(f"{archive_base}/{acc}-index.htm")
            or f"{archive_base}/{acc}-index.htm"
        )
        doc_type = "annual_report" if period_class == "annual" else "interim_report"
        year_match = re.search(r"20[12]\d", date or "")
        year = int(year_match.group()) if year_match else 0
        return {
            "url":          doc_url,
            "filename":     f"SEC_{form_clean}_{date}.html",
            "type":         doc_type,
            "priority":     _PRIORITY.get(form_clean, 2),
            "year":         year,
            "text":         f"{form_clean} filed {date}",
            "format":       "html",
            "source":       f"SEC EDGAR {form_clean}",
            "period_class": period_class,
        }

    def _fetch_and_process(recent_data: dict) -> None:
        form_types = recent_data.get("form",            [])
        accessions = recent_data.get("accessionNumber", [])
        dates      = recent_data.get("reportDate",      [])
        prim_docs  = recent_data.get("primaryDocument", [])
        for form, acc, date, prim_doc in zip(form_types, accessions, dates, prim_docs):
            form_clean = form.strip().upper()
            if len(annuals) >= max_annuals and len(quarterly_latest) >= max_quarterly:
                break
            if form_clean in _ANNUAL_FORMS and len(annuals) < max_annuals:
                annuals.append(_build_entry(form_clean, acc, date, prim_doc, "annual"))
            elif form_clean in _QUARTERLY_FORMS and len(quarterly_latest) < max_quarterly:
                quarterly_latest.append(_build_entry(form_clean, acc, date, prim_doc, "quarterly"))

    try:
        sub_url = f"https://data.sec.gov/submissions/CIK{cik_padded}.json"
        r = requests.get(sub_url, headers=_SEC_HEADERS, timeout=15)
        r.raise_for_status()
        time.sleep(0.15)
        payload = r.json()
        recent  = payload.get("filings", {}).get("recent", {})
        _fetch_and_process(recent)

        # If recent window didn't provide 3 annuals, try older filings pages
        if len(annuals) < max_annuals:
            for older_file in payload.get("filings", {}).get("files", []):
                if len(annuals) >= max_annuals:
                    break
                try:
                    older_url = f"https://data.sec.gov/submissions/{older_file['name']}"
                    ro = requests.get(older_url, headers=_SEC_HEADERS, timeout=15)
                    ro.raise_for_status()
                    time.sleep(0.15)
                    _fetch_and_process(ro.json())
                except Exception:
                    pass

    except Exception as exc:
        print(f"      SEC EDGAR submissions API Fehler ({ticker}): {exc}")

    result = annuals + quarterly_latest
    print(
        f"      SEC EDGAR: {len(annuals)} Jahresberichte + "
        f"{len(quarterly_latest)} Quartalsbericht(e) fuer {ticker} (CIK {cik_int})."
    )
    return result


# keep old name as alias so any existing callers continue to work
def find_sec_filings(ticker: str) -> list[dict]:
    """Deprecated alias — resolves CIK then calls get_sec_filings()."""
    cik = get_sec_cik(ticker)
    if not cik:
        return []
    return get_sec_filings(ticker, cik)


# ── 2b. HTML document loader ──────────────────────────────────────────────────

_HTML_PRIORITY_KEYWORDS = [
    "revenue", "net sales", "guidance", "outlook", "forward",
    "earnings per share", "eps", "free cash flow", "fcf",
    "operating income", "ebitda", "segment results", "fiscal year",
    "next year", "net income", "dividend", "buyback", "capex",
]

_MAX_HTML_CHARS = 60_000


def load_html_document(url: str, doc_type: str = "unknown",
                       ticker: str = "") -> list:
    """
    Fetches *url*, strips boilerplate HTML, extracts relevant text sections,
    caps at _MAX_HTML_CHARS (60 000 chars), and returns LangChain Document chunks.

    Priority sections (paragraphs/tables containing _HTML_PRIORITY_KEYWORDS)
    are prepended so they land in the first chunks and score highest in
    similarity search.

    Cache: ./ir_cache/{ticker}/{url_hash}.html  (24-hour TTL).
    SEC EDGAR URLs automatically use the SEC-compliant User-Agent.
    """
    url_hash   = hashlib.md5(url.encode()).hexdigest()[:8]
    # Per-ticker cache dir when ticker is known, otherwise shared _html dir
    cache_dir  = Path(CACHE_DIR) / (ticker if ticker else "_html")
    cache_path = cache_dir / f"{url_hash}.html"

    raw_text: str = ""
    if cache_path.exists():
        try:
            age_h = (time.time() - cache_path.stat().st_mtime) / 3600
            if age_h < CACHE_MAX_AGE_HOURS:
                raw_text = cache_path.read_text(encoding="utf-8")
        except Exception:
            pass

    if not raw_text:
        # SEC EDGAR requires a different User-Agent (email address mandatory)
        req_headers = _SEC_HEADERS if "sec.gov" in url else _HEADERS
        try:
            r = requests.get(url, headers=req_headers, timeout=30, allow_redirects=True)
            r.raise_for_status()
        except Exception as exc:
            print(f"      HTML laden fehlgeschlagen ({url[:60]}...): {exc}")
            return []

        soup = BeautifulSoup(r.text, "html.parser")

        # Remove noise tags
        for tag in soup(["script", "style", "nav", "header", "footer",
                         "aside", "noscript", "iframe", "form", "button"]):
            tag.decompose()

        # Prefer semantic content containers
        content_root = (
            soup.find("main")
            or soup.find("article")
            or soup.find(id=re.compile(r"content|main|body", re.I))
            or soup.body
            or soup
        )

        # Collect text blocks (paragraphs, table rows, headings)
        priority_blocks: list[str] = []
        normal_blocks:   list[str] = []

        for elem in content_root.find_all(
            ["p", "li", "td", "th", "h1", "h2", "h3", "h4", "caption"]
        ):
            text = elem.get_text(" ", strip=True)
            if len(text) < 20:
                continue
            lower = text.lower()
            if any(kw in lower for kw in _HTML_PRIORITY_KEYWORDS):
                priority_blocks.append(text)
            else:
                normal_blocks.append(text)

        combined = "\n".join(priority_blocks) + "\n" + "\n".join(normal_blocks)
        raw_text = combined[:_MAX_HTML_CHARS]

        # Persist to cache
        try:
            cache_dir.mkdir(parents=True, exist_ok=True)
            cache_path.write_text(raw_text, encoding="utf-8")
        except Exception:
            pass

    if not raw_text.strip():
        return []

    doc = Document(
        page_content=raw_text,
        metadata={"source": url, "type": doc_type, "format": "html", "page": "N/A"},
    )
    chunks = _splitter.split_documents([doc])
    for i, chunk in enumerate(chunks):
        chunk.metadata.update(
            {"source": url, "type": doc_type, "format": "html", "page": i}
        )
    return chunks


# ── 2. Find IR PDFs ───────────────────────────────────────────────────────────

def find_ir_pdfs(ir_url: str, ticker: str = "",
                  max_annual: int = 3, max_quarterly: int = 4,
                  wanted_years: set[int] | None = None) -> list[dict]:
    """
    Scrapes *ir_url* for IR documents.  Three-stage strategy:

    Stage 1 (always): scan for <a href="*.pdf"> links (existing logic).
    Stage 2 (PDF fallback): if no PDFs found, scan <a> text for IR keywords
      and return matching HTML links (marked format="html").
    Stage 3 (SEC fallback): if still empty AND ticker is a plain US ticker,
      call find_sec_filings() as last resort.

    *max_annual* / *max_quarterly* cap how many annual and interim documents
    are kept when *wanted_years* is not given. If *wanted_years* is given,
    annual docs are selected to cover exactly those fiscal years instead
    (see _deduplicate_and_spread).

    Returns each entry as:
      {url, filename, type, priority, year, text, format, source}
    """
    # ── Foreign Private Issuer: route directly to SEC EDGAR ──────────────────
    if ticker and ticker in _SEC_FOREIGN_FILERS:
        sec_ticker = _SEC_FOREIGN_FILERS[ticker]
        cik = get_sec_cik(sec_ticker)
        if cik:
            print(f"      {ticker} = Foreign Private Issuer, SEC EDGAR ({sec_ticker}, CIK {cik})")
            return get_sec_filings(ticker, cik)

    # ── SEC EDGAR URL: extract CIK and fetch filings directly ────────────────
    if ir_url and "sec.gov" in ir_url:
        cik_match = re.search(r"CIK=(\d+)", ir_url)
        cik = cik_match.group(1) if cik_match else None
        if not cik and ticker:
            cik = get_sec_cik(ticker)
        if cik:
            return get_sec_filings(ticker, cik)
        # CIK unresolvable — fall through to generic scraping below

    if not ir_url:
        if ticker and _sec_is_us_ticker(ticker):
            cik = get_sec_cik(ticker)
            if cik:
                return get_sec_filings(ticker, cik)
        return []

    parsed   = urlparse(ir_url)
    base_dom = f"{parsed.scheme}://{parsed.netloc}"

    # Try fetching with requests first; on 403/error fall through to Playwright below
    _requests_ok = False
    try:
        r = requests.get(ir_url, headers=_HEADERS, timeout=15)
        r.raise_for_status()
        _requests_ok = True
    except Exception:
        if ticker and _sec_is_us_ticker(ticker):
            cik = get_sec_cik(ticker)
            if cik:
                return get_sec_filings(ticker, cik)
        # Non-US: don't give up — Playwright (Stage 1c) will handle the 403

    soup = BeautifulSoup(r.text if _requests_ok else "", "html.parser")

    found: list[dict] = []
    seen:  set[str]   = set()

    def _resolve(href: str) -> str:
        href = href.strip()
        if href.startswith("//"):
            return parsed.scheme + ":" + href
        if href.startswith("/"):
            return base_dom + href
        if not href.startswith("http"):
            return urljoin(ir_url, href)
        return href

    def _is_pdf(href: str) -> bool:
        """True if the URL path (ignoring query string / fragment) ends with .pdf."""
        return href.split("?")[0].split("#")[0].lower().endswith(".pdf")

    # ── Stage 1: PDF links ────────────────────────────────────────────────────
    for a in soup.find_all("a", href=True):
        href = _resolve(a["href"])
        if not _is_pdf(href):
            continue
        if href in seen:
            continue
        seen.add(href)

        filename  = href.split("/")[-1].split("?")[0] or "document.pdf"
        link_text = a.get_text(" ", strip=True)
        label     = (link_text + " " + filename + " " + href).lower()

        if any(kw in label for kw in EXCLUDE_KEYWORDS):
            continue

        pdf_type = "other"
        priority = 99
        for t, p, any_kws, all_kws in _PDF_TYPE_RULES:
            hit_any = any(kw in label for kw in any_kws) if any_kws else True
            hit_all = all(kw in label for kw in all_kws) if all_kws else True
            if hit_any and hit_all:
                pdf_type = t
                priority = p
                break

        if pdf_type == "other":
            continue

        year_match = re.search(r"20[12][0-9]", label)
        year = int(year_match.group()) if year_match else 0

        found.append({
            "url":      href,
            "filename": filename,
            "type":     pdf_type,
            "priority": priority,
            "year":     year,
            "text":     link_text,
            "format":   "pdf",
            "source":   "IR website PDF",
        })

    if found:
        return _deduplicate_and_spread(found, max_annual=max_annual, max_latest=max_quarterly,
                                        wanted_years=wanted_years)

    # ── Stage 1b: Follow report sub-pages (EU/CH multi-level IR sites) ────────
    # Many European IR sites have PDFs only 1-2 clicks below the landing page.
    # We collect same-domain links whose path contains a report-index keyword,
    # visit each sub-page, and search for PDFs there.
    if not (ticker and _sec_is_us_ticker(ticker)):
        subpage_candidates: list[str] = []
        for a in soup.find_all("a", href=True):
            href = _resolve(a["href"])
            if urlparse(href).netloc != parsed.netloc:
                continue
            if href in seen or href == ir_url:
                continue
            path_lower = urlparse(href).path.lower().rstrip("/")
            if any(kw in path_lower for kw in _SUBPAGE_FOLLOW_PATTERNS):
                subpage_candidates.append(href)

        seen_sub: set[str] = set()
        for subpage_url in subpage_candidates[:8]:
            if subpage_url in seen_sub:
                continue
            seen_sub.add(subpage_url)
            print(f"      Suche PDFs auf Sub-Seite: {subpage_url[:80]}...")
            try:
                r2 = requests.get(subpage_url, headers=_HEADERS, timeout=12)
                if r2.status_code != 200 or len(r2.content) < 500:
                    continue
                soup2 = BeautifulSoup(r2.text, "html.parser")

                for a in soup2.find_all("a", href=True):
                    raw_href = a["href"].strip()
                    # Resolve relative to the sub-page URL (not the original ir_url)
                    if raw_href.startswith("//"):
                        href2 = parsed.scheme + ":" + raw_href
                    elif raw_href.startswith("/"):
                        href2 = base_dom + raw_href
                    elif not raw_href.startswith("http"):
                        href2 = urljoin(subpage_url, raw_href)
                    else:
                        href2 = raw_href

                    if not _is_pdf(href2):
                        continue
                    if href2 in seen:
                        continue
                    seen.add(href2)

                    filename  = href2.split("/")[-1].split("?")[0] or "document.pdf"
                    link_text = a.get_text(" ", strip=True)
                    label     = (link_text + " " + filename + " " + href2).lower()

                    if any(kw in label for kw in EXCLUDE_KEYWORDS):
                        continue

                    pdf_type, priority = "other", 99
                    for t, p, any_kws, all_kws in _PDF_TYPE_RULES:
                        hit_any = any(kw in label for kw in any_kws) if any_kws else True
                        hit_all = all(kw in label for kw in all_kws) if all_kws else True
                        if hit_any and hit_all:
                            pdf_type, priority = t, p
                            break

                    # EU/CH fallback: classify by filename pattern
                    if pdf_type == "other":
                        if _EU_PDF_ANNUAL.search(label):
                            pdf_type, priority = "annual_report", 3
                        elif _EU_PDF_INTERIM.search(label):
                            pdf_type, priority = "interim_report", 4
                        else:
                            continue

                    year_match = re.search(r"20[12][0-9]", label)
                    year = int(year_match.group()) if year_match else 0

                    found.append({
                        "url":      href2,
                        "filename": filename,
                        "type":     pdf_type,
                        "priority": priority,
                        "year":     year,
                        "text":     link_text,
                        "format":   "pdf",
                        "source":   "IR sub-page PDF",
                    })

            except Exception:
                continue

            if found:
                break  # stop after first sub-page that yields PDFs

        if found:
            return _deduplicate_and_spread(found, max_annual=max_annual, max_latest=max_quarterly,
                                        wanted_years=wanted_years)

    # ── Stage 1c: Playwright JS rendering (EU/CH sites with dynamic PDF lists) ──
    if not found and not (ticker and _sec_is_us_ticker(ticker)) and _PLAYWRIGHT_AVAILABLE:
        print(f"      Playwright-Rendering (multi-year): {ir_url[:60]}...")
        rendered = _playwright_render_multi_year(ir_url)
        if rendered:
            soup_pw = BeautifulSoup(rendered, "html.parser")

            def _scan_html_for_pdfs(html_src: str, soup_obj, source_label: str) -> None:
                """Extract PDF URLs from both <a href> and raw HTML regex (covers CDN embeds)."""
                # <a href> scan
                for a in soup_obj.find_all("a", href=True):
                    href = _resolve(a["href"])
                    if not _is_pdf(href) or href in seen:
                        continue
                    seen.add(href)
                    filename  = href.split("/")[-1].split("?")[0] or "document.pdf"
                    link_text = a.get_text(" ", strip=True)
                    label     = (link_text + " " + filename + " " + href).lower()
                    if any(kw in label for kw in EXCLUDE_KEYWORDS):
                        continue
                    pdf_type, priority = "other", 99
                    for t, p, any_kws, all_kws in _PDF_TYPE_RULES:
                        hit_any = any(kw in label for kw in any_kws) if any_kws else True
                        hit_all = all(kw in label for kw in all_kws) if all_kws else True
                        if hit_any and hit_all:
                            pdf_type, priority = t, p
                            break
                    if pdf_type == "other":
                        if _EU_PDF_ANNUAL.search(label):
                            pdf_type, priority = "annual_report", 3
                        elif _EU_PDF_INTERIM.search(label):
                            pdf_type, priority = "interim_report", 4
                        else:
                            continue
                    year_match = re.search(r"20[12][0-9]", label)
                    year = int(year_match.group()) if year_match else 0
                    found.append({
                        "url": href, "filename": filename, "type": pdf_type,
                        "priority": priority, "year": year, "text": link_text,
                        "format": "pdf", "source": source_label,
                    })
                # Regex scan — catches PDF URLs in JSON blobs / data-attrs not in <a href>
                for raw_url in re.findall(r'https?://[^\s"\'<>()\[\]]+\.pdf', html_src):
                    if raw_url in seen:
                        continue
                    seen.add(raw_url)
                    filename = raw_url.split("/")[-1].split("?")[0] or "document.pdf"
                    label    = (filename + " " + raw_url).lower()
                    if any(kw in label for kw in EXCLUDE_KEYWORDS):
                        continue
                    pdf_type, priority = "other", 99
                    for t, p, any_kws, all_kws in _PDF_TYPE_RULES:
                        hit_any = any(kw in label for kw in any_kws) if any_kws else True
                        hit_all = all(kw in label for kw in all_kws) if all_kws else True
                        if hit_any and hit_all:
                            pdf_type, priority = t, p
                            break
                    if pdf_type == "other":
                        if _EU_PDF_ANNUAL.search(label):
                            pdf_type, priority = "annual_report", 3
                        elif _EU_PDF_INTERIM.search(label):
                            pdf_type, priority = "interim_report", 4
                        else:
                            continue
                    year_match = re.search(r"20[12][0-9]", label)
                    year = int(year_match.group()) if year_match else 0
                    found.append({
                        "url": raw_url, "filename": filename, "type": pdf_type,
                        "priority": priority, "year": year, "text": "",
                        "format": "pdf", "source": source_label + " (regex)",
                    })

            # PDF links in JS-rendered page (<a href> + regex)
            _scan_html_for_pdfs(rendered, soup_pw, "Playwright IR page PDF")

            if found:
                return _deduplicate_and_spread(found, max_annual=max_annual, max_latest=max_quarterly,
                                        wanted_years=wanted_years)

            # Playwright sub-page following (for multi-level EU IR sites)
            seen_sub_pw: set[str] = set()
            for a in soup_pw.find_all("a", href=True):
                href = _resolve(a["href"])
                if urlparse(href).netloc != parsed.netloc:
                    continue
                if href in seen or href == ir_url or href in seen_sub_pw:
                    continue
                path_lower = urlparse(href).path.lower().rstrip("/")
                if not any(kw in path_lower for kw in _SUBPAGE_FOLLOW_PATTERNS):
                    continue
                seen_sub_pw.add(href)

            for subpage_url in list(seen_sub_pw)[:4]:
                print(f"      Playwright Sub-Seite: {subpage_url[:80]}...")
                rendered2 = _playwright_render(subpage_url)
                if not rendered2:
                    continue
                soup2 = BeautifulSoup(rendered2, "html.parser")
                _scan_html_for_pdfs(rendered2, soup2, "Playwright sub-page PDF")
                if found:
                    break

            if found:
                return _deduplicate_and_spread(found, max_annual=max_annual, max_latest=max_quarterly,
                                        wanted_years=wanted_years)

    # ── Stage 2: HTML links with IR keyword anchors ───────────────────────────
    print(f"      Keine PDFs gefunden auf {ir_url[:60]} — suche HTML IR-Links...")
    html_found: list[dict] = []

    for a in soup.find_all("a", href=True):
        href = _resolve(a["href"])
        if href in seen:
            continue
        # Skip obvious non-content URLs
        ext = href.split("?")[0].rsplit(".", 1)[-1].lower()
        if ext in ("jpg", "jpeg", "png", "gif", "svg", "css", "js", "xml", "zip"):
            continue
        seen.add(href)

        link_text = a.get_text(" ", strip=True)
        label     = (link_text + " " + href).lower()

        if any(kw in label for kw in EXCLUDE_KEYWORDS):
            continue
        if not any(kw in label for kw in _HTML_IR_KEYWORDS):
            continue

        # Classify using the same priority rules (applied to label)
        html_type = "other"
        priority  = 99
        for t, p, any_kws, all_kws in _PDF_TYPE_RULES:
            hit_any = any(kw in label for kw in any_kws) if any_kws else True
            hit_all = all(kw in label for kw in all_kws) if all_kws else True
            if hit_any and hit_all:
                html_type = t
                priority  = p
                break

        # Also classify by HTML IR keywords when PDF rules don't match
        if html_type == "other":
            if any(kw in label for kw in ["10-k", "10k", "annual report"]):
                html_type, priority = "annual_report", 3
            elif any(kw in label for kw in ["10-q", "10q", "quarterly results", "interim"]):
                html_type, priority = "interim_report", 4
            elif any(kw in label for kw in ["earnings", "press release", "financial results"]):
                html_type, priority = "earnings_release", 1
            elif any(kw in label for kw in ["presentation", "investor presentation"]):
                html_type, priority = "analyst_presentation", 2
            else:
                continue

        year_match = re.search(r"20[12][0-9]", label)
        year = int(year_match.group()) if year_match else 0
        filename = href.split("/")[-1].split("?")[0] or "page.html"
        if not filename.endswith((".htm", ".html", ".aspx", ".php")):
            filename += ".html"

        html_found.append({
            "url":      href,
            "filename": filename,
            "type":     html_type,
            "priority": priority,
            "year":     year,
            "text":     link_text,
            "format":   "html",
            "source":   "IR website HTML",
        })

    if html_found:
        return _deduplicate_and_spread(html_found, max_annual=max_annual, max_latest=max_quarterly,
                                        wanted_years=wanted_years)

    # ── Stage 3: SEC EDGAR fallback for US tickers ───────────────────────────
    if ticker and _sec_is_us_ticker(ticker):
        print(f"      Kein IR-Dokument auf Website — versuche SEC EDGAR fuer {ticker}...")
        cik = get_sec_cik(ticker)
        if cik:
            return get_sec_filings(ticker, cik)

    return []


# ── 3. Download PDF ───────────────────────────────────────────────────────────

def download_pdf(pdf_url: str, save_path: str) -> bool:
    """
    Downloads *pdf_url* to *save_path*.
    Enforces 50 MB size limit.  Returns True on success.
    """
    filename = os.path.basename(save_path)
    try:
        r = requests.get(pdf_url, headers=_HEADERS, timeout=30,
                         stream=True, allow_redirects=True)
        r.raise_for_status()

        # Pre-flight size check from Content-Length header
        content_length = int(r.headers.get("content-length", 0))
        if content_length > _MAX_FILE_BYTES:
            size_mb = content_length / 1024 / 1024
            print(f"      Überspringe {filename} ({size_mb:.1f}MB — Limit: 50MB)")
            return False

        os.makedirs(os.path.dirname(save_path), exist_ok=True)
        bytes_written = 0
        with open(save_path, "wb") as fh:
            for chunk in r.iter_content(chunk_size=65536):
                bytes_written += len(chunk)
                if bytes_written > _MAX_FILE_BYTES:
                    fh.close()
                    os.remove(save_path)
                    print(f"      Überspringe {filename} (>50MB während Download)")
                    return False
                fh.write(chunk)

        size_mb = bytes_written / 1024 / 1024
        print(f"      Lade PDF: {filename} ({size_mb:.1f}MB)")
        return True
    except Exception as exc:
        print(f"      Fehler beim Download {filename}: {exc}")
        return False


# ── 4. Load document ──────────────────────────────────────────────────────────

def load_document(file_path: str, doc_type: str = "unknown") -> list:
    """
    Loads *file_path* (.pdf / .pptx / .docx), splits into chunks.
    Returns list of Document chunks with metadata, or [] on failure.
    """
    filename = os.path.basename(file_path)
    ext = os.path.splitext(file_path)[1].lower()
    raw_docs: list[Document] = []

    try:
        if ext == ".pdf":
            loader   = PyPDFLoader(file_path)
            raw_docs = loader.load()

        elif ext == ".pptx":
            if not _PPTX_AVAILABLE:
                print(f"      python-pptx nicht installiert — überspringe {filename}")
                return []
            prs = _PptxPresentation(file_path)
            for slide_num, slide in enumerate(prs.slides, 1):
                parts: list[str] = []
                for shape in slide.shapes:
                    try:
                        txt = shape.text.strip()
                        if txt:
                            parts.append(txt)
                    except AttributeError:
                        pass
                try:
                    if slide.has_notes_slide:
                        notes = slide.notes_slide.notes_text_frame.text.strip()
                        if notes:
                            parts.append(f"[Notes] {notes}")
                except Exception:
                    pass
                slide_text = "\n".join(parts)
                if slide_text.strip():
                    raw_docs.append(Document(
                        page_content=slide_text,
                        metadata={"page": slide_num},
                    ))

        elif ext == ".docx":
            if not _DOCX_AVAILABLE:
                print(f"      Docx2txtLoader nicht verfügbar — überspringe {filename}")
                return []
            loader   = Docx2txtLoader(file_path)
            raw_docs = loader.load()

        else:
            return []

    except Exception as exc:
        print(f"      Fehler beim Laden {filename}: {exc}")
        return []

    chunks = _splitter.split_documents(raw_docs)
    for chunk in chunks:
        chunk.metadata.update({
            "source":   filename,
            "type":     doc_type,
            "page":     chunk.metadata.get("page", "N/A"),
        })
    return chunks


# ── 5. Build vectorstore ──────────────────────────────────────────────────────

def build_vectorstore(ticker: str, documents: list, source_urls: list[str]):
    """
    Returns FAISS vectorstore.
    Loads from cache if metadata.json is younger than CACHE_MAX_AGE_HOURS;
    otherwise builds from *documents* and persists to cache.
    Returns None if *documents* is empty and no valid cache exists.
    """
    cache_dir = Path(CACHE_DIR) / ticker
    meta_file = cache_dir / "metadata.json"
    index_file = cache_dir / "index.faiss"

    # Try loading valid cache
    if meta_file.exists() and index_file.exists():
        try:
            meta  = json.loads(meta_file.read_text(encoding="utf-8"))
            age_h = (time.time() - meta["timestamp"]) / 3600
            if age_h < CACHE_MAX_AGE_HOURS:
                vs = FAISS.load_local(str(cache_dir), _get_emb(),
                                      allow_dangerous_deserialization=True)
                return vs
        except Exception:
            pass  # corrupt / stale — rebuild

    if not documents:
        return None

    vs = FAISS.from_documents(documents, _get_emb())

    cache_dir.mkdir(parents=True, exist_ok=True)
    vs.save_local(str(cache_dir))
    meta_file.write_text(json.dumps({
        "ticker":      ticker,
        "timestamp":   time.time(),
        "source_urls": source_urls,
        "doc_count":   len(documents),
    }), encoding="utf-8")

    return vs


# ── 6. @tool get_ir_analysis ──────────────────────────────────────────────────

_SYNTHESIS_SYSTEM = """\
You are a financial analyst. Extract structured data from IR document excerpts.
Return ONLY valid JSON — no explanatory text before or after.
Do NOT invent numbers. If a figure is not found in the context, write "not found".
Always cite the page or section for every extracted figure.\
"""

_SYNTHESIS_HUMAN = """\
Du bist ein Senior Buy-Side Analyst. Extrahiere präzise Finanzdaten für {company} ({ticker}) aus den vorliegenden IR-Auszügen.

WICHTIGSTE REGELN FÜR EINHEITEN:
- Alle Felder, die auf '_bn' enden, müssen in MILLIARDEN (Billions) angegeben werden.
- Beispiel: Wenn im Bericht '420.1 Millionen CHF' steht, schreibe 0.4201.
- Beispiel: Wenn im Bericht '1.2 Milliarden CHF' steht, schreibe 1.2.
- Rechne Millionen-Beträge konsequent durch 1000.

HIERARCHIE DER DATEN:
1. Suche zuerst nach der Tabelle 'Key Figures', 'Financial Highlights' oder 'Group Overview'.
2. Bevorzuge 'Adjusted' (bereinigte) Werte gegenüber 'Reported' Werten.
3. Bevorzuge 'Restated' (angepasste) Vorjahreswerte, falls vorhanden (wichtig bei Kapitalerhöhungen).

yfinance Referenzdaten (nutze diese NUR zum Abgleich, nicht zum Kopieren):
{yf_data}

Extrahiere die Daten und antworte AUSSCHLIESSLICH mit diesem JSON-Schema:
{{
  "revenue_bn": <float or "not found">,
  "revenue_currency": "<CHF, USD, EUR>",
  "ebitda_bn": <float - absolut aus Erfolgsrechnung extrahieren oder berechnen>,
  "ebitda_margin_pct": <float>,
  "ebit_bn": <float - 'Operating Result' oder 'EBIT' nutzen>,
  "recurring_ebit_margin_pct": <float - oft als 'EBIT margin before special items' bezeichnet>,
  "net_income_bn": <float>,
  "adjusted_eps": <float - Gewinn pro Aktie bereinigt>,
  "adjusted_eps_note": "<Quelle und Grund der Bereinigung>",
  "free_cashflow_bn": <float>,
  "free_cashflow_note": "<Quelle>",
  "net_debt_bn": <float - positiv=Schulden, negativ=Netto-Cash>,
  "total_equity_bn": <float - aus der Bilanz>,
  "total_assets_bn": <float - Bilanzsumme>,
  "tax_rate_pct": <float - effektiver Steuersatz für ROIC>,
  "invested_capital_bn": <float - falls explizit genannt>,
  "dividend_per_share": <float>,
  "dividend_currency": "<CHF, USD, EUR>",
  "guidance_2026": "<Zitat zur Guidance für Umsatz/Marge>",
  "guidance_2027": "<Zitat zur Guidance>",
  "consensus_eps_2026": <float - falls 'Consensus' Tabelle im IR gefunden wurde>,
  "consensus_eps_2027": <float>,
  "consensus_revenue_2026_bn": <float>,
  "consensus_revenue_2027_bn": <float>,
  "management_tone": "<optimistic, neutral, cautious>",
  "key_statements": ["<Direktzitat inkl. Seite [P. X]>"],
  "pe_distortion_explanation": "<Grund für KGV-Verzerrung (z.B. Kapitalerhöhung 2025, Restructuring)>",
  "data_quality": "<high, medium, low>",
  "yfinance_discrepancies": ["<Beschreibung, wenn IR-Wert >10% von yf_data abweicht>"]
}}

IR document excerpts:
{context}

yfinance data for cross-checking:
{yf_data}

Return ONLY this JSON structure:
{{
  "adjusted_eps": <float or "not found">,
  "adjusted_eps_note": "<source page/section or explanation>",
  "free_cashflow_bn": <float in billions of home currency or "not found">,
  "free_cashflow_currency": "<CHF, USD, EUR, GBP or 'not found'>",
  "free_cashflow_note": "<source page/section>",
  "revenue_bn": <float in billions or "not found">,
  "revenue_currency": "<CHF, USD, EUR, GBP or 'not found'>",
  "ebitda_margin_pct": <float or "not found">,
  "recurring_ebit_margin_pct": <float or "not found">,
  "net_debt_bn": <float or "not found">,
  "capex_bn": <float or "not found">,
  "dividend_per_share": <float or "not found">,
  "dividend_currency": "<CHF, USD, EUR, GBP or 'not found'>",
  "guidance_2026": "<exact quote from document or 'not found'>",
  "guidance_2027": "<exact quote from document or 'not found'>",
  "consensus_eps_2026": <float or "not found">,
  "consensus_eps_2027": <float or "not found">,
  "consensus_eps_2028": <float or "not found">,
  "consensus_revenue_2026_bn": <float or "not found">,
  "consensus_revenue_2027_bn": <float or "not found">,
  "consensus_revenue_2028_bn": <float or "not found">,
  "management_tone": "<optimistic, neutral, or cautious>",
  "key_statements": ["<direct quote [Page X]>", "<direct quote [Page X]>"],
  "pe_distortion_explanation": "<explanation if P/E distorted by spin-off or one-time item, else 'none'>",
  "ir_sources": [],
  "data_quality": "<high, medium, or low>",
  "yfinance_discrepancies": ["<description of any figure differing from yfinance by >10%>"]
}}\
"""

_ANNUAL_SYSTEM = """\
You are a financial analyst. Extract structured annual financial data from IR document excerpts.
Return ONLY valid JSON — no text before or after.
Do NOT invent numbers. Write "not found" for any figure not present in the context.
Always note the fiscal year for each entry.\
"""

_ANNUAL_HUMAN = """\
Du bist ein Senior Buy-Side Analyst. Extrahiere Jahresdaten für {company} ({ticker}) aus den vorliegenden Jahresbericht-Auszügen.

EINHEITEN: Alle '_bn'-Felder in MILLIARDEN. Millionen durch 1000 teilen.
HIERARCHIE: 1) Key Figures / Financial Highlights  2) Adjusted-Werte  3) Restated-Vorjahre

Extrahiere Daten für JEDES im Context enthaltene Geschäftsjahr und gib eine Liste zurück (neuestes Jahr zuerst, max. 10 Jahre).
Guidance und Consensus NUR aus den jüngsten Zahlen (letzter Jahresbericht).

Antworte NUR mit diesem JSON:
{{
  "years": [
    {{
      "fiscal_year": <int z.B. 2024>,
      "period_end": "<echtes Ende des Geschäftsjahres als YYYY-MM-DD, z.B. bei \
vom Kalenderjahr abweichendem GJ (z.B. 'Ende Januar'/'Ende September' im Bericht) \
- or 'not found' wenn nicht im Dokument angegeben>",
      "revenue_bn": <float or "not found">,
      "revenue_currency": "<CHF/USD/EUR>",
      "ebitda_bn": <float or "not found">,
      "ebitda_margin_pct": <float or "not found">,
      "ebit_bn": <float or "not found">,
      "recurring_ebit_margin_pct": <float or "not found">,
      "net_income_bn": <float or "not found">,
      "adjusted_eps": <float or "not found">,
      "free_cashflow_bn": <float or "not found">,
      "net_debt_bn": <float positiv=Schulden, negativ=Netto-Cash, or "not found">,
      "dividend_per_share": <float or "not found">,
      "data_quality": "<high/medium/low>"
    }}
  ],
  "guidance_2026": "<Zitat or 'not found'>",
  "guidance_2027": "<Zitat or 'not found'>",
  "consensus_eps_2026": <float or "not found">,
  "consensus_eps_2027": <float or "not found">,
  "consensus_eps_2028": <float or "not found">,
  "consensus_revenue_2026_bn": <float or "not found">,
  "consensus_revenue_2027_bn": <float or "not found">,
  "consensus_revenue_2028_bn": <float or "not found">,
  "management_tone": "<optimistic/neutral/cautious>",
  "key_statements": ["<Direktzitat inkl. Seite [P. X]>"],
  "pe_distortion_explanation": "<Grund für KGV-Verzerrung or 'none'>",
  "data_quality": "<high/medium/low>",
  "yfinance_discrepancies": ["<Beschreibung wenn IR-Wert >10% von yf_data abweicht>"]
}}

IR-Auszüge (Jahresberichte):
{context}

yfinance Referenzdaten (NUR zum Abgleich):
{yf_data}
"""

_QUARTERLY_MULTI_SYSTEM = """\
You are a financial analyst. Extract structured interim/quarterly financial data \
from IR document excerpts covering possibly several interim periods of the current \
fiscal year (e.g. Q1, H1, 9M). Return ONLY valid JSON — no text before or after.
Do NOT invent numbers. Write "not found" for any figure not present in the context.
Always note the fiscal year and quarter label for each entry.\
"""

_QUARTERLY_MULTI_HUMAN = """\
Du bist ein Senior Buy-Side Analyst. Extrahiere Zwischenbericht-Daten für {company} ({ticker}) \
aus den vorliegenden Quartals-/Interim-Auszügen des laufenden Geschäftsjahres.

EINHEITEN: Alle '_bn'-Felder in MILLIARDEN. Millionen durch 1000 teilen.
HIERARCHIE: 1) Key Figures / Financial Highlights  2) Adjusted-Werte  3) kumulierte (YTD) Werte, falls das Quartal selbst nicht separat ausgewiesen ist.

Der Context kann Auszüge aus MEHREREN Zwischenberichten enthalten (z.B. Q1 UND 9M desselben Jahres).
Extrahiere für JEDEN im Context erkennbaren, distinkten Berichtszeitraum einen eigenen Eintrag \
(neuestes Quartal zuerst). Gib KEINE Guidance- oder Consensus-Werte an (nur aus Jahresberichten).

Antworte NUR mit diesem JSON:
{{
  "periods": [
    {{
      "fiscal_year": <int z.B. 2026>,
      "quarter": "<z.B. Q1 2026, H1 2026, 9M 2026>",
      "period_end": "<YYYY-MM-DD or 'not found'>",
      "revenue_bn": <float or "not found">,
      "revenue_currency": "<CHF/USD/EUR>",
      "ebitda_bn": <float or "not found">,
      "ebitda_margin_pct": <float or "not found">,
      "ebit_bn": <float or "not found">,
      "net_income_bn": <float or "not found">,
      "adjusted_eps": <float or "not found">,
      "free_cashflow_bn": <float or "not found">,
      "net_debt_bn": <float positiv=Schulden, negativ=Netto-Cash, or "not found">,
      "yoy_comparable_growth_pct": <float or "not found">,
      "data_quality": "<high/medium/low>"
    }}
  ]
}}

IR-Auszüge (Quartals-/Zwischenberichte):
{context}

yfinance Referenzdaten (NUR zum Abgleich):
{yf_data}
"""

_EMPTY_IR: dict = {
    "adjusted_eps":              "not found",
    "adjusted_eps_note":         "not found",
    "free_cashflow_bn":          "not found",
    "free_cashflow_currency":    "not found",
    "free_cashflow_note":        "not found",
    "revenue_bn":                "not found",
    "revenue_currency":          "not found",
    "ebitda_margin_pct":         "not found",
    "recurring_ebit_margin_pct": "not found",
    "net_debt_bn":               "not found",
    "capex_bn":                  "not found",
    "dividend_per_share":        "not found",
    "dividend_currency":         "not found",
    "guidance_2026":             "not found",
    "guidance_2027":             "not found",
    "consensus_eps_2026":        "not found",
    "consensus_eps_2027":        "not found",
    "consensus_eps_2028":        "not found",
    "consensus_revenue_2026_bn": "not found",
    "consensus_revenue_2027_bn": "not found",
    "consensus_revenue_2028_bn": "not found",
    "management_tone":           "neutral",
    "key_statements":            [],
    "pe_distortion_explanation": "none",
    "ir_sources":                [],
    "data_quality":              "low",
    "yfinance_discrepancies":    [],
    # Multi-year extension (new)
    "ir_annual_years":           [],   # list of dicts, one per fiscal year
    "ir_quarterly_latest":       None, # dict with latest quarterly metrics, or None
    "ir_quarterly_periods":      [],   # list of dicts, one per interim period this fiscal year
}


@tool
def get_ir_analysis(ticker: str) -> dict:
    """
    RAG-basierte Analyse von IR-Dokumenten (Geschäftsbericht, Consensus Estimates,
    Earnings Presentation). Extrahiert EPS, FCF, Guidance und Konsensdaten direkt
    aus den Quelldokumenten. Liefert strukturiertes JSON inkl. yfinance-Crosscheck.
    """
    # Get company name from yfinance
    company_name = ticker
    try:
        company_name = yf.Ticker(ticker).info.get("longName", ticker)
    except Exception:
        pass

    cache_dir = Path(CACHE_DIR) / ticker
    meta_file = cache_dir / "metadata.json"

    # ── Check whether a fresh cache already covers us ────────────────────────
    cache_fresh  = False
    source_urls: list[str] = []

    if meta_file.exists():
        try:
            meta  = json.loads(meta_file.read_text(encoding="utf-8"))
            age_h = (time.time() - meta["timestamp"]) / 3600
            if age_h < CACHE_MAX_AGE_HOURS:
                cache_fresh = True
                source_urls = meta.get("source_urls", [])
        except Exception:
            pass

    # ── Steps 1-4: find, download, process (skipped when cache is fresh) ─────
    all_chunks: list = []
    wanted_years: set[int] | None = None

    if not cache_fresh:
        # Step 1: Find IR URL
        ir_url = find_ir_url(ticker, company_name)
        if not ir_url:
            print(f"      IR-Seite fuer {ticker} nicht gefunden.")

        # Step 2: Find documents (PDFs, HTML, or SEC filings).
        # Gap-driven annual depth: only fetch/extract fiscal years the DB is
        # actually missing (up to a 10-year window), plus the latest year
        # (to catch restatements) — a mature ticker tops up just the gap
        # instead of re-extracting years already cached. Quarterly/interim
        # docs are always fetched in full, since that's where
        # forward-estimate-relevant news is.
        from datetime import datetime as _dt_now
        current_yr = _dt_now.now().year
        target_years = set(range(current_yr - 10, current_yr))  # last 10 completed FYs
        try:
            from tools.financial_db import get_annual_years_present
            present_years = get_annual_years_present(ticker)
        except Exception:
            present_years = set()
        wanted_years = (target_years - present_years) | {current_yr - 1}
        pdfs = find_ir_pdfs(ir_url, ticker=ticker, max_quarterly=4, wanted_years=wanted_years)
        if not pdfs:
            print(f"      Keine IR-Dokumente fuer {ticker} gefunden.")

        # Step 3 + 4: Load each document — route by format
        for doc_info in pdfs:
            fmt          = doc_info.get("format", "pdf")
            period_class = doc_info.get("period_class", "annual")
            fiscal_year  = doc_info.get("year", 0)

            if fmt == "html":
                print(f"      Lade HTML ({period_class}, {fiscal_year}): {doc_info['url'][:60]}...")
                chunks = load_html_document(
                    doc_info["url"], doc_type=doc_info["type"], ticker=ticker
                )
                if chunks:
                    source_label = doc_info.get("source", "IR website HTML")
                    for chunk in chunks:
                        chunk.metadata["source"]       = source_label
                        chunk.metadata["period_class"] = period_class
                        chunk.metadata["fiscal_year"]  = fiscal_year
                    all_chunks.extend(chunks)
                    source_urls.append(doc_info["url"])
            else:
                save_path = str(cache_dir / doc_info["filename"])
                if download_pdf(doc_info["url"], save_path):
                    chunks = load_document(save_path, doc_type=doc_info["type"])
                    if chunks:
                        for chunk in chunks:
                            chunk.metadata["period_class"] = period_class
                            chunk.metadata["fiscal_year"]  = fiscal_year
                        all_chunks.extend(chunks)
                        source_urls.append(doc_info["url"])

    if not all_chunks and not cache_fresh:
        result = {**_EMPTY_IR, "error": f"Keine IR-Dokumente für {ticker} gefunden."}
        result["ir_sources"] = source_urls
        return result

    # ── Step 5: Build / load vectorstore ────────────────────────────────────
    vs = build_vectorstore(ticker, all_chunks, source_urls)
    if vs is None:
        result = {**_EMPTY_IR, "error": "Vectorstore konnte nicht erstellt werden."}
        result["ir_sources"] = source_urls
        return result

    # yfinance cashflow for cross-check (shared by both passes)
    yf_data: dict = {}
    try:
        yf_data = get_cashflow_data.invoke(ticker)
    except Exception:
        pass

    def _build_context(filter_val: str, char_cap: int = 10000, k: int = 4) -> str:
        """Retrieve chunks filtered by period_class and build context string."""
        parts: list[str] = []
        for query in STANDARD_QUERIES:
            try:
                # Use metadata filter; fetch_k ensures enough candidates across years
                hits = vs.similarity_search(
                    query, k=k,
                    filter={"period_class": filter_val},
                    fetch_k=max(20, k * 5),
                )
                if not hits:
                    # Fallback: unfiltered if the filter returns nothing
                    hits = vs.similarity_search(query, k=2)
                if hits:
                    parts.append(f"\n=== {query.upper()} ===")
                    for doc in hits:
                        page = doc.metadata.get("page", "N/A")
                        src  = doc.metadata.get("source", "N/A")
                        yr   = doc.metadata.get("fiscal_year", "")
                        parts.append(
                            f"[Page {page} | {src} | year={yr}] "
                            f"{doc.page_content[:400]}"
                        )
            except Exception:
                pass
        return "\n".join(parts)[:char_cap]

    def _safe_parse(raw: str) -> dict | None:
        try:
            s = raw.find("{")
            e = raw.rfind("}") + 1
            if s != -1 and e > 0:
                return json.loads(raw[s:e])
        except Exception:
            pass
        return None

    # ── Pass A: Annual extraction (gap-driven, up to 10 years) ───────────────
    annual_k = min(max(len(wanted_years) if wanted_years else 3, 3) + 2, 12)
    annual_context = _build_context("annual", char_cap=20000, k=annual_k)

    annual_prompt = ChatPromptTemplate.from_messages([
        ("system", _ANNUAL_SYSTEM),
        ("human",  _ANNUAL_HUMAN),
    ])
    ir_annual_years: list[dict] = []
    annual_top_level: dict = {}
    try:
        raw_annual = (annual_prompt | _get_llm() | StrOutputParser()).invoke({
            "ticker":  ticker,
            "company": company_name,
            "context": annual_context,
            "yf_data": json.dumps(yf_data, ensure_ascii=False),
        })
        parsed_annual = _safe_parse(raw_annual) or {}
        ir_annual_years = parsed_annual.get("years", [])
        # Promote top-level fields from annual extraction
        annual_top_level = {k: v for k, v in parsed_annual.items() if k != "years"}
    except Exception as exc:
        print(f"      IR Annual-Extraktion fehlgeschlagen ({ticker}): {exc}")

    # ── Pass B: Quarterly extraction (all interim periods published this fiscal year) ──
    ir_quarterly_periods: list[dict] = []
    ir_quarterly_latest: dict | None = None
    has_quarterly_chunks = any(
        c.metadata.get("period_class") == "quarterly" for c in (all_chunks or [])
    )
    if has_quarterly_chunks:
        quarterly_context = _build_context("quarterly", char_cap=8000)
        quarterly_prompt = ChatPromptTemplate.from_messages([
            ("system", _QUARTERLY_MULTI_SYSTEM),
            ("human",  _QUARTERLY_MULTI_HUMAN),
        ])
        try:
            raw_q = (quarterly_prompt | _get_llm() | StrOutputParser()).invoke({
                "ticker":  ticker,
                "company": company_name,
                "context": quarterly_context,
                "yf_data": json.dumps(yf_data, ensure_ascii=False),
            })
            parsed_q = _safe_parse(raw_q) or {}
            ir_quarterly_periods = parsed_q.get("periods", [])
        except Exception as exc:
            print(f"      IR Quarterly-Extraktion fehlgeschlagen ({ticker}): {exc}")

        if ir_quarterly_periods:
            ir_quarterly_latest = ir_quarterly_periods[0]
            print(f"      IR: {len(ir_quarterly_periods)} Zwischenberichte extrahiert "
                  f"({', '.join(str(p.get('quarter', '?')) for p in ir_quarterly_periods)})")

    # ── Assemble result ───────────────────────────────────────────────────────
    # Top-level single-year fields from most recent annual (backward compat)
    if ir_annual_years:
        latest_annual = ir_annual_years[0]
        result = {**_EMPTY_IR}
        for field in (
            "revenue_bn", "revenue_currency", "ebitda_bn", "ebitda_margin_pct",
            "ebit_bn", "recurring_ebit_margin_pct", "net_income_bn", "adjusted_eps",
            "free_cashflow_bn", "net_debt_bn", "dividend_per_share",
        ):
            if latest_annual.get(field) not in (None, "not found"):
                result[field] = latest_annual[field]
        # Guidance/consensus/tone from annual top-level
        result.update({k: v for k, v in annual_top_level.items()
                       if k in _EMPTY_IR and k not in result})
    else:
        # Fallback: single-pass synthesis with existing prompt
        fallback_context = _build_context("annual", char_cap=12000)
        fallback_prompt = ChatPromptTemplate.from_messages([
            ("system", _SYNTHESIS_SYSTEM),
            ("human",  _SYNTHESIS_HUMAN),
        ])
        try:
            raw_fb = (fallback_prompt | _get_llm() | StrOutputParser()).invoke({
                "ticker":  ticker,
                "company": company_name,
                "context": fallback_context,
                "yf_data": json.dumps(yf_data, ensure_ascii=False),
            })
            result = _safe_parse(raw_fb) or {**_EMPTY_IR}
        except Exception as exc:
            result = {**_EMPTY_IR, "error": str(exc)}

    result["ir_annual_years"]      = ir_annual_years
    result["ir_quarterly_latest"]  = ir_quarterly_latest
    result["ir_quarterly_periods"] = ir_quarterly_periods
    result["ir_sources"]           = source_urls
    return result


# ── 7. Consensus estimates from IR ───────────────────────────────────────────

_DERIVE_SYSTEM = """\
Du bist ein erfahrener Buy-Side Analyst der Forward-Estimates erstellt.
Nutze IR-Guidance, historische Wachstumsraten, Management Tone und Sektordynamiken.
Sei konservativ und erkläre jede Annahme explizit mit ihrer Herleitung.
Antworte AUSSCHLIESSLICH mit validem JSON.\
"""

_GUIDANCE_DERIVE_HUMAN = """\
Unternehmen: {ticker} (Sektor: {sector})

MANAGEMENT GUIDANCE (aus IR-Dokument):
Guidance 2026: {guidance_2026}
Guidance 2027: {guidance_2027}

AKTUELLE KENNZAHLEN (letztes Geschäftsjahr):
Umsatz: {revenue_bn} Mrd. {currency}
EBITDA-Marge: {ebitda_margin_pct}%
EPS (bereinigt): {adjusted_eps}

HISTORISCHE DATEN:
{historical_data}

AUFGABE:
1. Parse den Guidance-Text und extrahiere konkrete Zielwerte für {year_e1}
2. Extrapoliere {year_e2} und {year_e3} mit historischem CAGR und Sektordurchschnitt
3. Sei konservativ - im Zweifelsfall niedrigere Wachstumsrate

Antworte NUR mit diesem JSON:
{{
  "estimates": {{
    "{year_e1}": {{"revenue_bn": <Mrd.>, "ebitda_margin_pct": <Prozent>, "eps": <EPS>, "source": "Management Guidance"}},
    "{year_e2}": {{"revenue_bn": <Mrd.>, "ebitda_margin_pct": <Prozent>, "eps": <EPS>, "source": "extrapoliert"}},
    "{year_e3}": {{"revenue_bn": <Mrd.>, "ebitda_margin_pct": <Prozent>, "eps": <EPS>, "source": "extrapoliert"}}
  }},
  "key_assumptions": [
    "<Annahme 1: konkrete Guidance-Quelle und Herleitung>",
    "<Annahme 2: CAGR-Annahme für {year_e2}/{year_e3} mit Begründung>"
  ]
}}\
"""

_DERIVE_HUMAN = """\
Erstelle 3-Jahres Forward-Estimates für {ticker} (Sektor: {sector}).

IR-ANALYSE:
{ir_output}

HISTORISCHE DATEN:
{historical_data}

Konfidenz-Regeln:
- hoch:    IR-Guidance vorhanden UND historische Daten vollständig
- mittel:  nur historische Daten ODER unvollständige Guidance
- niedrig: keine Guidance, nur Sektordurchschnitte als Basis

JSON-Format (kein erklärender Text):
{{
  "source": "LLM-Ableitung aus IR-Dokumenten und historischen Daten",
  "confidence": "hoch" | "mittel" | "niedrig",
  "estimates": {{
    "{year_e1}": {{"revenue_bn": <Mrd.>, "ebitda_margin_pct": <Prozent>, "eps": <EPS>, "source": "derived"}},
    "{year_e2}": {{"revenue_bn": <Mrd.>, "ebitda_margin_pct": <Prozent>, "eps": <EPS>, "source": "derived"}},
    "{year_e3}": {{"revenue_bn": <Mrd.>, "ebitda_margin_pct": <Prozent>, "eps": <EPS>, "source": "derived"}}
  }},
  "key_assumptions": [
    "<Annahme 1: Herleitung + Quelle>",
    "<Annahme 2: Herleitung + Quelle>",
    "<Annahme 3: Herleitung + Quelle>"
  ],
  "disclaimer": "Diese Schätzungen stammen aus öffentlichen IR-Dokumenten / wurden vom LLM abgeleitet. Schätzung aus IR-Daten."
}}\
"""

def _make_estimate_fallback(base_year: int) -> dict:
    e1, e2, e3 = f"{base_year+1}E", f"{base_year+2}E", f"{base_year+3}E"
    return {
        "source":      "LLM-Ableitung — Analyse nicht verfügbar",
        "confidence":  "niedrig",
        "estimates": {
            e1: {"revenue_bn": "-", "ebitda_margin_pct": "-", "eps": "-", "source": "-"},
            e2: {"revenue_bn": "-", "ebitda_margin_pct": "-", "eps": "-", "source": "-"},
            e3: {"revenue_bn": "-", "ebitda_margin_pct": "-", "eps": "-", "source": "-"},
        },
        "key_assumptions": [],
        "disclaimer": (
            "Diese Schätzungen stammen aus öffentlichen IR-Dokumenten / wurden vom LLM abgeleitet. "
            "Schätzung aus IR-Daten."
        ),
    }


def _base_year_from_historical(historical_data: dict) -> int:
    """Ermittelt das letzte abgeschlossene Geschäftsjahr aus historischen Daten."""
    import re as _re
    year_keys = sorted(
        k for k in historical_data
        if _re.match(r"20[12]\d[AE]?$", str(k))
    )
    if year_keys:
        return int(_re.sub(r"[AE]", "", year_keys[-1]))
    from datetime import date
    return date.today().year - 1


def consensus_estimates_from_ir(
    ticker: str,
    ir_output: dict,
    historical_data: dict,
    sector: str = "",
) -> dict:
    """
    Returns 3-year forward estimates dynamically based on the last completed fiscal year.
    NOT a @tool - called directly by fundamental_agent.

    Fast path: if the IR document already contains consensus_eps and
    consensus_revenue, those are returned directly without an LLM call.
    """
    def _is_found(v) -> bool:
        return v not in (None, "not found", "n/v", "-", "")

    base_year = _base_year_from_historical(historical_data)
    e1 = f"{base_year + 1}E"
    e2 = f"{base_year + 2}E"
    e3 = f"{base_year + 3}E"
    y1, y2, y3 = base_year + 1, base_year + 2, base_year + 3

    # ── Fast path: direct IR consensus figures ───────────────────────────────
    eps_e1 = ir_output.get(f"consensus_eps_{y1}")
    eps_e2 = ir_output.get(f"consensus_eps_{y2}")
    eps_e3 = ir_output.get(f"consensus_eps_{y3}")
    rev_e1 = ir_output.get(f"consensus_revenue_{y1}_bn")
    rev_e2 = ir_output.get(f"consensus_revenue_{y2}_bn")
    rev_e3 = ir_output.get(f"consensus_revenue_{y3}_bn")

    if all(_is_found(v) for v in [eps_e1, eps_e2, eps_e3, rev_e1, rev_e2, rev_e3]):
        return {
            "source":     "IR-Dokument (Consensus Estimates)",
            "confidence": "hoch",
            "estimates": {
                e1: {"revenue_bn": rev_e1, "ebitda_margin_pct": ir_output.get("ebitda_margin_pct", "n/v"), "eps": eps_e1, "source": "direct from IR"},
                e2: {"revenue_bn": rev_e2, "ebitda_margin_pct": "-", "eps": eps_e2, "source": "direct from IR"},
                e3: {"revenue_bn": rev_e3, "ebitda_margin_pct": "-", "eps": eps_e3, "source": "direct from IR"},
            },
            "key_assumptions": [
                f"Direkt aus IR-Consensus-Dokument entnommen ({', '.join(ir_output.get('ir_sources', ['IR']))}).",
                "Keine LLM-Ableitung — Zahlen stammen unverändert aus dem Unternehmensdokument.",
            ],
            "methodology": "Zahlen direkt aus dem vom Unternehmen publizierten Consensus Sheet übernommen.",
            "disclaimer": (
                "Quelle: IR Consensus Sheet (direkt vom Unternehmen). "
                "Kein Ersatz für professionelle Konsensdaten."
            ),
        }

    # ── Priority 2: Management Guidance from IR ───────────────────────────────
    guidance_e1 = ir_output.get(f"guidance_{y1}")
    guidance_e2 = ir_output.get(f"guidance_{y2}")
    revenue_bn    = ir_output.get("revenue_bn")
    ebitda_margin = ir_output.get("ebitda_margin_pct")

    if _is_found(guidance_e1) and _is_found(revenue_bn) and _is_found(ebitda_margin):
        guidance_prompt = ChatPromptTemplate.from_messages([
            ("system", _DERIVE_SYSTEM),
            ("human",  _GUIDANCE_DERIVE_HUMAN),
        ])
        try:
            raw_json: str = (guidance_prompt | _get_llm() | StrOutputParser()).invoke({
                "ticker":            ticker,
                "sector":            sector or "unbekannt",
                "guidance_2026":     guidance_e1,
                "guidance_2027":     guidance_e2 if _is_found(guidance_e2) else "nicht verfügbar",
                "revenue_bn":        revenue_bn,
                "currency":          ir_output.get("revenue_currency", ""),
                "ebitda_margin_pct": ebitda_margin,
                "adjusted_eps":      ir_output.get("adjusted_eps", "nicht verfügbar"),
                "historical_data":   json.dumps(historical_data, ensure_ascii=False),
                "year_e1": e1, "year_e2": e2, "year_e3": e3,
            })
            s = raw_json.find("{")
            e_pos = raw_json.rfind("}") + 1
            if s != -1 and e_pos > 0:
                parsed = json.loads(raw_json[s:e_pos])
                ir_sources = ir_output.get("ir_sources", ["IR-Dokument"])
                return {
                    "source":     "Management Guidance (IR-Dokument)",
                    "confidence": "mittel-hoch",
                    "estimates":  parsed["estimates"],
                    "key_assumptions": parsed.get("key_assumptions", [])
                        + [f"Guidance-Quelle: {', '.join(ir_sources)}"],
                    "methodology": (
                        f"{e1} direkt aus Management Guidance abgeleitet. "
                        f"{e2}/{e3} via historischem CAGR extrapoliert."
                    ),
                    "disclaimer": (
                        "Quelle: Management Guidance (IR-Dokument). "
                        "Kein Ersatz für professionelle Konsensdaten."
                    ),
                }
        except Exception:
            pass  # fall through to Priority 3

    # ── Priority 3: LLM derivation fallback ──────────────────────────────────
    prompt = ChatPromptTemplate.from_messages([
        ("system", _DERIVE_SYSTEM),
        ("human",  _DERIVE_HUMAN),
    ])

    try:
        raw_json: str = (prompt | _get_llm() | StrOutputParser()).invoke({
            "ticker":          ticker,
            "sector":          sector or "unbekannt",
            "ir_output":       json.dumps(ir_output,       ensure_ascii=False),
            "historical_data": json.dumps(historical_data, ensure_ascii=False),
            "year_e1": e1, "year_e2": e2, "year_e3": e3,
        })
        s = raw_json.find("{")
        e_pos = raw_json.rfind("}") + 1
        if s != -1 and e_pos > 0:
            result = json.loads(raw_json[s:e_pos])
            result.setdefault("methodology", "LLM-Ableitung aus IR-Dokumenten und historischen Daten.")
            result["disclaimer"] = (
                "Quelle: LLM-Ableitung aus IR-Dokumenten. "
                "Kein Ersatz für professionelle Konsensdaten."
            )
            return result
    except Exception:
        pass

    return _make_estimate_fallback(base_year)


# ── Smoke test ────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    import json
    from dotenv import load_dotenv
    load_dotenv()

    print("=== IR RAG Tool Test: HOLN.SW ===")
    result = get_ir_analysis.invoke("HOLN.SW")
    print(json.dumps(result, indent=2, ensure_ascii=False, default=str))